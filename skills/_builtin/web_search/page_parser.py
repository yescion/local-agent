"""Multi-format web page / API response parsing."""

from __future__ import annotations

import csv
import io
import json
import re
import xml.etree.ElementTree as ET
import zipfile
from html import unescape
from pathlib import Path
from urllib.parse import urlparse

try:
    from charset_normalizer import from_bytes

    _HAS_CHARSET_NORMALIZER = True
except ImportError:
    _HAS_CHARSET_NORMALIZER = False

try:
    from trafilatura import extract as trafilatura_extract
    from trafilatura import html2txt as trafilatura_html2txt

    _HAS_TRAFILATURA = True
except ImportError:
    _HAS_TRAFILATURA = False

_HTML_MIMES = frozenset(
    {
        "text/html",
        "application/xhtml+xml",
    }
)
_JSON_MIMES = frozenset(
    {
        "application/json",
        "application/ld+json",
        "text/json",
        "application/vnd.api+json",
    }
)
_XML_MIMES = frozenset(
    {
        "application/xml",
        "text/xml",
        "application/rss+xml",
        "application/atom+xml",
        "application/rdf+xml",
    }
)
_YAML_MIMES = frozenset(
    {
        "application/x-yaml",
        "application/yaml",
        "text/yaml",
        "text/x-yaml",
    }
)
_CSV_MIMES = frozenset({"text/csv", "application/csv"})
_MARKDOWN_MIMES = frozenset({"text/markdown", "text/x-markdown"})
_PLAIN_MIMES = frozenset({"text/plain"})
_RTF_MIMES = frozenset({"text/rtf", "application/rtf"})
_CODE_MIMES = frozenset(
    {
        "text/javascript",
        "application/javascript",
        "text/css",
        "application/typescript",
    }
)

_BINARY_MIME_MAP: dict[str, str] = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "doc",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xls",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/epub+zip": "epub",
}

_EXT_FORMAT_MAP: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "doc",
    ".xlsx": "xlsx",
    ".xls": "xls",
    ".pptx": "pptx",
    ".epub": "epub",
    ".json": "json",
    ".xml": "xml",
    ".yaml": "yaml",
    ".yml": "yaml",
    ".csv": "csv",
    ".md": "markdown",
    ".markdown": "markdown",
    ".rtf": "rtf",
}


def normalize_mime(content_type: str) -> str:
    return content_type.split(";", 1)[0].strip().lower()


def decode_bytes(content: bytes, content_type: str) -> str:
    charset = _charset_from_header(content_type)
    if charset:
        try:
            return content.decode(charset, errors="replace")
        except LookupError:
            pass

    head = content[:8192].decode("ascii", errors="ignore")
    match = re.search(
        r'<meta[^>]+charset=["\']?([\w-]+)',
        head,
        re.IGNORECASE,
    )
    if match:
        try:
            return content.decode(match.group(1), errors="replace")
        except LookupError:
            pass

    if _HAS_CHARSET_NORMALIZER:
        result = from_bytes(content).best()
        if result is not None:
            return str(result)

    return content.decode("utf-8", errors="replace")


def parse_response(content: bytes, content_type: str, *, url: str = "") -> tuple[str, str, str]:
    """Return (title, body, format_label) from raw HTTP response bytes."""
    mime = normalize_mime(content_type)
    binary_fmt = _resolve_binary_format(content, mime, url)
    if binary_fmt:
        return _parse_binary(content, binary_fmt)

    text = decode_bytes(content, content_type)
    return parse_text(text, mime, url=url, raw=content)


def parse_text(
    text: str,
    mime: str,
    *,
    url: str = "",
    raw: bytes | None = None,
) -> tuple[str, str, str]:
    """Return (title, body, format_label) from decoded text."""
    mime = normalize_mime(mime)
    stripped = text.lstrip("\ufeff").strip()

    if raw is not None and mime in ("application/octet-stream", "binary/octet-stream", ""):
        binary_fmt = _resolve_binary_format(raw, mime, url)
        if binary_fmt:
            return _parse_binary(raw, binary_fmt)

    if mime in _HTML_MIMES or _looks_like_html(stripped):
        return _parse_html(stripped, url=url)
    if mime in _JSON_MIMES or _looks_like_json(stripped):
        return _parse_json(stripped)
    if mime in _XML_MIMES or _looks_like_xml(stripped):
        return _parse_xml(stripped)
    if mime in _YAML_MIMES or _looks_like_yaml(stripped):
        return _parse_yaml(stripped)
    if mime in _CSV_MIMES:
        return _parse_csv(stripped)
    if mime in _MARKDOWN_MIMES:
        return _parse_markdown(stripped)
    if mime in _RTF_MIMES or stripped.startswith("{\\rtf"):
        return _parse_rtf(stripped)
    if mime in _PLAIN_MIMES:
        return "", stripped, "text"
    if mime in _CODE_MIMES:
        return "", stripped, mime.split("/")[-1]
    if mime.startswith("text/"):
        if mime in _MARKDOWN_MIMES:
            return _parse_markdown(stripped)
        return "", stripped, "text"

    sniffed = _sniff_format(stripped, url)
    if sniffed == "html":
        return _parse_html(stripped, url=url)
    if sniffed == "json":
        return _parse_json(stripped)
    if sniffed == "xml":
        return _parse_xml(stripped)
    if sniffed == "yaml":
        return _parse_yaml(stripped)
    if sniffed == "markdown":
        return _parse_markdown(stripped)
    if sniffed == "csv":
        return _parse_csv(stripped)
    if sniffed == "rtf":
        return _parse_rtf(stripped)
    if stripped:
        return "", stripped, "text"

    return "", "", "unknown"


def _resolve_binary_format(content: bytes, mime: str, url: str) -> str | None:
    if mime in _BINARY_MIME_MAP:
        return _BINARY_MIME_MAP[mime]

    magic_fmt = _sniff_binary_magic(content)
    if magic_fmt:
        return magic_fmt

    ext = Path(urlparse(url).path).suffix.lower()
    ext_fmt = _EXT_FORMAT_MAP.get(ext)
    if ext_fmt and ext_fmt in {"pdf", "docx", "xlsx", "pptx", "epub", "doc", "xls"}:
        verified = _verify_binary_format(content, ext_fmt)
        if verified:
            return verified
    return None


def _verify_binary_format(content: bytes, fmt: str) -> str | None:
    magic = _sniff_binary_magic(content)
    if magic == fmt:
        return fmt
    if fmt == "docx" and magic == "docx":
        return "docx"
    if fmt in {"xlsx", "xls"} and magic in {"xlsx", "xls"}:
        return magic
    if fmt == "pdf" and content.startswith(b"%PDF"):
        return "pdf"
    if fmt == "epub" and magic == "epub":
        return "epub"
    return None


def _sniff_binary_magic(content: bytes) -> str | None:
    if content.startswith(b"%PDF"):
        return "pdf"
    if content.startswith(b"{\\rtf"):
        return "rtf"
    if content[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
        return "doc"
    if content.startswith(b"PK\x03\x04"):
        return _sniff_zip_format(content)
    return None


def _sniff_zip_format(content: bytes) -> str | None:
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            names = set(zf.namelist())
            if "word/document.xml" in names:
                return "docx"
            if "xl/workbook.xml" in names:
                return "xlsx"
            if "ppt/presentation.xml" in names:
                return "pptx"
            if "META-INF/container.xml" in names:
                return "epub"
    except zipfile.BadZipFile:
        pass
    return None


def _parse_binary(content: bytes, fmt: str) -> tuple[str, str, str]:
    parsers = {
        "pdf": _parse_pdf,
        "docx": _parse_docx,
        "doc": _parse_legacy_office,
        "xlsx": _parse_xlsx,
        "xls": _parse_legacy_office,
        "pptx": _parse_pptx,
        "epub": _parse_epub,
        "rtf": lambda c: _parse_rtf(c.decode("utf-8", errors="replace")),
    }
    parser = parsers.get(fmt)
    if parser is None:
        return "", f"错误：不支持的二进制格式 — {fmt}", fmt
    return parser(content)


def _charset_from_header(content_type: str) -> str | None:
    match = re.search(r"charset=([\w-]+)", content_type, re.IGNORECASE)
    return match.group(1) if match else None


def _looks_like_html(text: str) -> bool:
    head = text[:4096].lower()
    return any(
        token in head
        for token in ("<html", "<body", "<!doctype html", "<article", "<head")
    )


def _looks_like_json(text: str) -> bool:
    return text[:1] in "{["


def _looks_like_xml(text: str) -> bool:
    return text[:1] == "<" and not _looks_like_html(text)


def _looks_like_yaml(text: str) -> bool:
    head = text[:200].lstrip()
    if not head or head[0] in "{[":
        return False
    return bool(re.match(r"^[\w.-]+\s*:", head))


def _sniff_format(text: str, url: str) -> str:
    if not text:
        return "unknown"
    if _looks_like_json(text):
        return "json"
    if _looks_like_html(text):
        return "html"
    if text[:1] == "<":
        return "xml"
    if _looks_like_yaml(text):
        return "yaml"
    if text.startswith("{\\rtf"):
        return "rtf"
    ext = Path(urlparse(url).path).suffix.lower()
    if ext in {".md", ".markdown"}:
        return "markdown"
    if ext == ".csv" and "," in text[:500]:
        return "csv"
    return "text"


def _pick_main_html_fragment(html: str) -> str | None:
    stripped = re.sub(
        r"(?is)<(script|style|noscript|svg|iframe)[^>]*>.*?</\1>",
        " ",
        html,
    )
    for pattern in (
        r"(?is)<article[^>]*>.*?</article>",
        r"(?is)<main[^>]*>.*?</main>",
    ):
        match = re.search(pattern, stripped)
        if match and match.group(0).strip():
            return match.group(0)
    match = re.search(
        r'(?is)<div[^>]+(?:id|class)=["\'][^"\']*(?:article|content|post|entry|main)[^"\']*["\'][^>]*>.*?</div>',
        stripped,
    )
    if match and len(re.sub(r"<[^>]+>", "", match.group(0)).strip()) > 200:
        return match.group(0)
    return None


def _parse_html(html: str, *, url: str) -> tuple[str, str, str]:
    title = _html_title(html)
    source = _pick_main_html_fragment(html) or html

    if _HAS_TRAFILATURA:
        body = trafilatura_extract(
            source,
            url=url or None,
            output_format="txt",
            include_tables=True,
            favor_precision=False,
        )
        if not body:
            body = trafilatura_html2txt(source)
        if body:
            return title, body.strip(), "html"

    _, body = _parse_html_fallback(html)
    return title, body, "html"


def _html_title(html: str) -> str:
    match = re.search(r"<title[^>]*>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
    return _clean_html(match.group(1)) if match else ""


def _parse_html_fallback(html: str) -> tuple[str, str]:
    title = _html_title(html)
    stripped = re.sub(r"(?is)<(script|style|noscript|svg|iframe)[^>]*>.*?</\1>", " ", html)
    for pattern in (
        r"(?is)<article[^>]*>(.*?)</article>",
        r"(?is)<main[^>]*>(.*?)</main>",
    ):
        match = re.search(pattern, stripped)
        if match and match.group(1).strip():
            return title, _html_fragment_to_text(match.group(1))
    body_match = re.search(r"(?is)<body[^>]*>(.*?)</body>", stripped)
    chunk = body_match.group(1) if body_match else stripped
    return title, _html_fragment_to_text(chunk)


def _html_fragment_to_text(html: str) -> str:
    html = re.sub(r"(?is)<br\s*/?>", "\n", html)
    html = re.sub(r"(?is)</(p|div|h[1-6]|li|tr|section|blockquote)>", "\n", html)
    text = _clean_html(html)
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _clean_html(text: str) -> str:
    text = re.sub(r"<[^>]+>", "", text)
    return unescape(text).strip()


def _parse_json(text: str) -> tuple[str, str, str]:
    data = json.loads(text)
    title = ""
    if isinstance(data, dict):
        for key in ("title", "name", "subject", "headline"):
            value = data.get(key)
            if isinstance(value, str) and value.strip():
                title = value.strip()
                break
    body = json.dumps(data, ensure_ascii=False, indent=2)
    return title, body, "json"


def _parse_yaml(text: str) -> tuple[str, str, str]:
    try:
        import yaml
    except ImportError:
        return "", "错误：未安装 pyyaml", "yaml"

    data = yaml.safe_load(text)
    title = ""
    if isinstance(data, dict):
        for key in ("title", "name", "subject"):
            value = data.get(key)
            if isinstance(value, str) and value.strip():
                title = value.strip()
                break
    body = yaml.dump(data, allow_unicode=True, sort_keys=False).strip()
    return title, body, "yaml"


def _parse_csv(text: str) -> tuple[str, str, str]:
    sample = text[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel

    reader = csv.reader(io.StringIO(text), dialect)
    rows = [row for row in reader if any(cell.strip() for cell in row)]
    if not rows:
        return "", "", "csv"

    lines: list[str] = []
    if len(rows) > 1:
        header = rows[0]
        lines.append(" | ".join(header))
        lines.append(" | ".join("---" for _ in header))
        for row in rows[1:]:
            lines.append(" | ".join(row))
    else:
        lines.append(" | ".join(rows[0]))
    return "", "\n".join(lines), "csv"


def _parse_markdown(text: str) -> tuple[str, str, str]:
    try:
        import frontmatter

        post = frontmatter.loads(text)
        title = ""
        if isinstance(post.metadata, dict):
            raw_title = post.metadata.get("title")
            if isinstance(raw_title, str) and raw_title.strip():
                title = raw_title.strip()
        body = (post.content or "").strip()
        return title, body or text.strip(), "markdown"
    except Exception:
        title_match = re.search(r"^#\s+(.+)$", text, re.MULTILINE)
        title = title_match.group(1).strip() if title_match else ""
        return title, text.strip(), "markdown"


def _parse_rtf(text: str) -> tuple[str, str, str]:
    body = re.sub(r"\\[a-z]+\d*\s?", "", text)
    body = re.sub(r"[{}]", "", body)
    body = re.sub(r"\s+", " ", body).strip()
    return "", body, "rtf"


def _parse_xml(text: str) -> tuple[str, str, str]:
    root = ET.fromstring(text)
    root_tag = _local_tag(root.tag)
    if root_tag in {"rss", "feed", "rdf"}:
        return _parse_feed(root)
    title = _element_text(_find_descendant(root, "title"))
    body = _xml_to_text(root)
    return title, body.strip(), "xml"


def _find_child(parent: ET.Element, local_name: str) -> ET.Element | None:
    for child in parent:
        if _local_tag(child.tag) == local_name:
            return child
    return None


def _find_descendant(parent: ET.Element, local_name: str) -> ET.Element | None:
    if _local_tag(parent.tag) == local_name:
        return parent
    for element in parent.iter():
        if element is not parent and _local_tag(element.tag) == local_name:
            return element
    return None


def _parse_feed(root: ET.Element) -> tuple[str, str, str]:
    channel = _find_descendant(root, "channel")
    if channel is None:
        channel = root
    title = _element_text(_find_child(channel, "title"))
    if not title and _local_tag(root.tag) == "feed":
        title = _element_text(_find_child(root, "title"))

    lines: list[str] = []
    if title:
        lines.append(f"Feed: {title}")

    for item in root.iter():
        tag = _local_tag(item.tag)
        if tag not in {"item", "entry"}:
            continue
        item_title = _element_text(_find_child(item, "title"))
        summary_el = _find_child(item, "description")
        if summary_el is None:
            summary_el = _find_child(item, "summary")
        if summary_el is None:
            summary_el = _find_child(item, "content")
        summary = _element_text(summary_el)
        link = _element_text(_find_child(item, "link"))
        if not link:
            link_el = _find_child(item, "link")
            if link_el is not None:
                link = (link_el.attrib.get("href") or "").strip()

        if item_title:
            lines.append(f"- {item_title}")
        if summary:
            lines.append(f"  {summary}")
        if link:
            lines.append(f"  {link}")

    body = "\n".join(lines).strip()
    return title, body, "feed"


def _parse_pdf(content: bytes) -> tuple[str, str, str]:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "", "错误：未安装 pypdf（pip install pypdf）", "pdf"

    reader = PdfReader(io.BytesIO(content))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text.strip())
    body = "\n\n".join(pages).strip()
    title = (reader.metadata.title or "").strip() if reader.metadata else ""
    if not body:
        return title, "未能从 PDF 提取文本（可能是扫描件）", "pdf"
    return title, body, "pdf"


def _parse_docx(content: bytes) -> tuple[str, str, str]:
    try:
        from docx import Document
    except ImportError:
        return "", "错误：未安装 python-docx（pip install python-docx）", "docx"

    doc = Document(io.BytesIO(content))
    title = ""
    lines: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = getattr(para.style, "name", "") or ""
        if not title and "heading" in style_name.lower():
            title = text
        lines.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))

    if not title and doc.core_properties.title:
        title = doc.core_properties.title.strip()
    body = "\n".join(lines).strip()
    if not body:
        return title, "未能从 DOCX 提取文本", "docx"
    return title, body, "docx"


def _parse_xlsx(content: bytes) -> tuple[str, str, str]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "", "错误：未安装 openpyxl（pip install openpyxl）", "xlsx"

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines: list[str] = []
    title = wb.sheetnames[0] if wb.sheetnames else ""
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if value is None else str(value) for value in row]
            if any(cell.strip() for cell in cells):
                lines.append(" | ".join(cells))
    wb.close()
    body = "\n".join(lines).strip()
    if not body:
        return title, "未能从 XLSX 提取数据", "xlsx"
    return title, body, "xlsx"


def _parse_pptx(content: bytes) -> tuple[str, str, str]:
    try:
        from pptx import Presentation
    except ImportError:
        return (
            "",
            "错误：未安装 python-pptx，无法解析 PPTX（pip install python-pptx）",
            "pptx",
        )

    prs = Presentation(io.BytesIO(content))
    title = ""
    lines: list[str] = []
    for slide_no, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                slide_lines.append(text)
        if slide_lines:
            if not title:
                title = slide_lines[0]
            lines.append(f"## Slide {slide_no}")
            lines.extend(slide_lines)
    body = "\n".join(lines).strip()
    if not body:
        return title, "未能从 PPTX 提取文本", "pptx"
    return title, body, "pptx"


def _parse_epub(content: bytes) -> tuple[str, str, str]:
    title = ""
    parts: list[str] = []
    with zipfile.ZipFile(io.BytesIO(content)) as zf:
        for name in zf.namelist():
            if name.endswith(".opf"):
                opf = zf.read(name).decode("utf-8", errors="replace")
                match = re.search(r"<dc:title[^>]*>(.*?)</dc:title>", opf, re.IGNORECASE | re.DOTALL)
                if match and not title:
                    title = _clean_html(match.group(1))

        for name in sorted(zf.namelist()):
            lower = name.lower()
            if not lower.endswith((".xhtml", ".html", ".htm")):
                continue
            if any(skip in lower for skip in ("nav.xhtml", "/nav/", "toc.")):
                continue
            html = zf.read(name).decode("utf-8", errors="replace")
            _, text, _ = _parse_html(html, url="")
            if text:
                parts.append(text)

    body = "\n\n".join(parts).strip()
    if not body:
        return title, "未能从 EPUB 提取文本", "epub"
    return title, body, "epub"


def _parse_legacy_office(content: bytes) -> tuple[str, str, str]:
    fmt = "doc" if content[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" else "xls"
    return (
        "",
        f"错误：旧版 Office 格式（.{fmt}）暂不支持，请使用 .docx/.xlsx 或 PDF",
        fmt,
    )


def _local_tag(tag: str) -> str:
    return tag.rsplit("}", 1)[-1].lower()


def _element_text(element: ET.Element | None) -> str:
    if element is None:
        return ""
    return "".join(element.itertext()).strip()


def _xml_to_text(element: ET.Element) -> str:
    parts: list[str] = []
    for child in element.iter():
        if child is element:
            continue
        text = (child.text or "").strip()
        if text:
            parts.append(text)
    if parts:
        return "\n".join(parts)
    return _element_text(element)
